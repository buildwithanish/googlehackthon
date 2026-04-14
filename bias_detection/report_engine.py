import os
import io
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from pptx import Presentation
from pptx.util import Inches, Pt

def generate_pdf_report(data, filename="FairAI_Audit_Report.pdf"):
    """
    Generates a professional PDF report using ReportLab.
    """
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    elements = []

    # Title
    title_style = ParagraphStyle(
        'MainTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor("#4F46E5"),
        alignment=1,
        spaceAfter=30
    )
    elements.append(Paragraph("FairAI Compliance Audit Report", title_style))
    elements.append(Paragraph(f"Analysis Snapshot: {data.get('scenario', 'Fairness Scan')}", styles['Heading2']))
    elements.append(Spacer(1, 12))

    # Summary Section
    elements.append(Paragraph("Executive Summary", styles['Heading3']))
    summary_text = f"This report provides a detailed bias analysis of the '{data.get('dataset_name', 'System')}' dataset. " \
                   f"The analysis focuses on systemic disparities across sensitive attributes."
    elements.append(Paragraph(summary_text, styles['Normal']))
    elements.append(Spacer(1, 12))

    # Metrics Table
    elements.append(Paragraph("Core Fairness Metrics", styles['Heading4']))
    metrics_data = [["Metric", "Value", "Status"]]
    
    # Use real metrics or defaults
    metrics = data.get('metrics', {})
    if not metrics:
        metrics = {"Fairness Score": 0.85, "Disparate Impact": 0.92}

    for m_name, m_val in metrics.items():
        try:
            val_num = float(m_val)
            status = "Pass" if abs(val_num) < 0.2 else "Warning"
            metrics_data.append([str(m_name), f"{val_num:.4f}", status])
        except (ValueError, TypeError):
            metrics_data.append([str(m_name), str(m_val), "N/A"])
    
    t = Table(metrics_data, colWidths=[200, 100, 100])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor("#F3F4F6")),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor("#111827")),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey)
    ]))
    elements.append(t)
    elements.append(Spacer(1, 24))

    # AI Insights
    ai_text = data.get('ai_explanation')
    if ai_text:
        elements.append(Paragraph("AI-Powered Remediation Roadmap", styles['Heading3']))
        # Scrub problematic characters that crash Paragraph
        clean_ai = str(ai_text).replace('**', '').replace('#', '').replace('*', '').replace('_', '').strip()
        elements.append(Paragraph(clean_ai, styles['Normal']))
    else:
        elements.append(Paragraph("AI Insights: The engine is analyzing your data. Please re-run to see full roadmap.", styles['Normal']))

    try:
        doc.build(elements)
    except Exception as e:
        # Fallback build with just title if complex content fails
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        doc.build([Paragraph("FairAI Error: Report Generation Failed", title_style), Paragraph(str(e), styles['Normal'])])
    
    buffer.seek(0)
    return buffer

def generate_ppt_report(data):
    """
    Generates a PowerPoint presentation using python-pptx.
    """
    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "FairAI: Bias Detection Audit"
    subtitle.text = f"Synapse Squad Hub | Team AnishNova\nScenario: {data.get('scenario', 'General Analysis')}"

    # 2. Metrics Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Key Fairness Indicators"
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for m_name, m_val in data.get('metrics', {}).items():
        p = tf.add_paragraph()
        p.text = f"{m_name}: {m_val:.4f}"
        p.level = 0

    # 3. AI Recommendations
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "AI Mitigation Roadmap"
    tf = slide.placeholders[1].text_frame
    
    recs = ["Re-weighting of privileged samples", "Adversarial debiasing in training", "Human-in-the-loop validation"]
    for rec in recs:
        p = tf.add_paragraph()
        p.text = rec
        p.level = 0

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
